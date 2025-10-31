#!/usr/bin/env python3
"""
Convert Quarto chapter QMD files to PowerPoint presentations.
Uses tex2img for LaTeX equation rendering and python-pptx for slide generation.
"""

import re
import os
import sys
import argparse
from pathlib import Path
import json
from typing import List, Dict, Tuple, Optional

# Check for required packages
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    import matplotlib
    matplotlib.use('Agg')  # Non-interactive backend
    import matplotlib.pyplot as plt
    from matplotlib import mathtext
except ImportError:
    print("Error: matplotlib not installed. Run: pip install matplotlib")
    sys.exit(1)


class ChapterParser:
    """Parse QMD chapter files and extract structure."""

    def __init__(self, qmd_file: str):
        self.qmd_file = Path(qmd_file)
        self.content = self.qmd_file.read_text(encoding='utf-8')
        self.macros = self._load_macros()

    def _load_macros(self) -> str:
        """Load LaTeX macros from macros.qmd if it exists."""
        macros_file = self.qmd_file.parent / 'macros.qmd'
        if macros_file.exists():
            content = macros_file.read_text(encoding='utf-8')
            # Extract newcommand definitions
            macro_pattern = r'\\newcommand\{[^}]+\}\{[^}]+\}'
            macros = re.findall(macro_pattern, content)
            return '\n'.join(macros)
        return ''

    def extract_title(self) -> str:
        """Extract chapter title from markdown."""
        match = re.search(r'^# (.+?)(?:\{|$)', self.content, re.MULTILINE)
        return match.group(1).strip() if match else "Untitled"

    def extract_sections(self) -> List[Dict]:
        """Extract major sections (## headers) with content."""
        sections = []

        # Find all ## sections
        pattern = r'^## (.+?)$'
        matches = list(re.finditer(pattern, self.content, re.MULTILINE))

        for i, match in enumerate(matches):
            title = match.group(1).strip()
            start = match.end()

            # Find content until next ## or end of file
            if i < len(matches) - 1:
                end = matches[i + 1].start()
            else:
                end = len(self.content)

            content = self.content[start:end].strip()

            sections.append({
                'title': title,
                'content': content,
                'subsections': self._extract_subsections(content)
            })

        return sections

    def _extract_subsections(self, content: str) -> List[Dict]:
        """Extract ### and #### subsections from content."""
        subsections = []

        # Match both ### and ####
        pattern = r'^(###|####) (.+?)$'
        matches = list(re.finditer(pattern, content, re.MULTILINE))

        for i, match in enumerate(matches):
            level = len(match.group(1))
            title = match.group(2).strip()
            start = match.end()

            if i < len(matches) - 1:
                end = matches[i + 1].start()
            else:
                end = len(content)

            sub_content = content[start:end].strip()

            subsections.append({
                'level': level,
                'title': title,
                'content': sub_content
            })

        return subsections

    def extract_equations(self) -> List[Dict]:
        """Extract all LaTeX equations (display and inline)."""
        equations = []

        # Display equations: $$...$$
        display_pattern = r'\$\$([^$]+?)\$\$'
        for match in re.finditer(display_pattern, self.content, re.DOTALL):
            equations.append({
                'type': 'display',
                'latex': match.group(1).strip(),
                'position': match.start()
            })

        # Inline equations: $...$
        # Be careful not to match $$ which we already got
        inline_pattern = r'(?<!\$)\$(?!\$)([^$\n]+?)\$(?!\$)'
        for match in re.finditer(inline_pattern, self.content):
            equations.append({
                'type': 'inline',
                'latex': match.group(1).strip(),
                'position': match.start()
            })

        # Sort by position
        equations.sort(key=lambda x: x['position'])

        return equations

    def extract_figures(self) -> List[Dict]:
        """Extract Python figure code blocks."""
        figures = []

        # Pattern for figure blocks with labels
        pattern = r'```\{python\}[^`]*?#\| label: (fig-[^\n]+)[^`]*?```'

        for match in re.finditer(pattern, self.content, re.DOTALL):
            code = match.group(0)
            label = match.group(1).strip()

            # Extract caption if present
            caption_match = re.search(r'#\| fig-cap: (.+?)$', code, re.MULTILINE)
            caption = caption_match.group(1).strip() if caption_match else ''

            figures.append({
                'label': label,
                'caption': caption,
                'code': code
            })

        return figures

    def extract_iframes(self) -> List[Dict]:
        """Extract iframe embeds."""
        iframes = []

        pattern = r'<iframe[^>]*src="([^"]+)"[^>]*>.*?</iframe>'

        for match in re.finditer(pattern, self.content, re.DOTALL):
            url = match.group(1)

            # Try to find associated caption
            # Look for text after iframe or figure label
            context_start = max(0, match.start() - 200)
            context_end = min(len(self.content), match.end() + 200)
            context = self.content[context_start:context_end]

            label_match = re.search(r'\{#(fig-[^\}]+)\}', context)
            label = label_match.group(1) if label_match else ''

            iframes.append({
                'url': url,
                'label': label
            })

        return iframes

    def extract_callouts(self, content: str) -> List[Dict]:
        """Extract callout boxes (Principle, Rule, etc.)."""
        callouts = []

        pattern = r':::\ (\w+)\s*$(.+?)^:::'

        for match in re.finditer(pattern, content, re.MULTILINE | re.DOTALL):
            callout_type = match.group(1)
            callout_content = match.group(2).strip()

            callouts.append({
                'type': callout_type,
                'content': callout_content
            })

        return callouts


class EquationRenderer:
    """Render LaTeX equations to PNG images using matplotlib."""

    def __init__(self, output_dir: Path, preamble: str = ''):
        self.output_dir = output_dir
        self.output_dir.mkdir(exist_ok=True)
        self.preamble = preamble

    def render_equation(self, latex: str, filename: str, mode: str = 'display', dpi: int = 300) -> Path:
        """Render a single equation to PNG using matplotlib."""
        try:
            fig = plt.figure(figsize=(10, 2))
            fig.patch.set_facecolor('white')
            fig.patch.set_alpha(0.0)  # Transparent background

            ax = fig.add_subplot(111)
            ax.axis('off')

            # LaTeX is already extracted without delimiters, wrap in $ or $$
            if mode == 'display':
                latex_wrapped = f'${latex}$'
                fontsize = 24
            else:
                latex_wrapped = f'${latex}$'
                fontsize = 18

            ax.text(0.5, 0.5, latex_wrapped,
                   fontsize=fontsize,
                   ha='center', va='center',
                   transform=ax.transAxes)

            output_path = self.output_dir / filename
            plt.savefig(str(output_path), dpi=dpi, bbox_inches='tight',
                       pad_inches=0.1, transparent=True)
            plt.close(fig)

            return output_path
        except Exception as e:
            # Silently skip equations that can't be rendered
            # print(f"Warning: Failed to render equation: {latex[:50]}...")
            plt.close('all')  # Clean up
            return None

    def render_all(self, equations: List[Dict]) -> Dict[int, Path]:
        """Render all equations and return mapping of index to file path."""
        rendered = {}

        for i, eq in enumerate(equations):
            filename = f"equation_{i:03d}.png"
            mode = 'inline' if eq['type'] == 'inline' else 'display'

            path = self.render_equation(eq['latex'], filename, mode=mode)
            if path:
                rendered[i] = path

        return rendered


class SlideBuilder:
    """Build PowerPoint presentation from parsed chapter."""

    # Color scheme (professional finance/academic theme)
    HEADER_COLOR = RGBColor(0, 51, 102)  # Navy blue
    ACCENT_COLOR = RGBColor(255, 107, 53)  # Orange
    TEXT_COLOR = RGBColor(0, 0, 0)  # Black

    def __init__(self, template: Optional[str] = None):
        if template and Path(template).exists():
            self.prs = Presentation(template)
        else:
            self.prs = Presentation()

    def add_title_slide(self, title: str, subtitle: str = ''):
        """Add title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])

        slide.shapes.title.text = title
        if subtitle and len(slide.shapes) > 1:
            slide.placeholders[1].text = subtitle

    def add_section_slide(self, title: str):
        """Add section divider slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])  # Section header layout
        slide.shapes.title.text = title

    def add_content_slide(self, title: str, bullets: List[str] = None,
                         equation_images: List[Path] = None):
        """Add content slide with bullets and/or equations."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # Title and content
        slide.shapes.title.text = title

        if bullets:
            # Add bullet points
            if len(slide.placeholders) > 1:
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()

                for bullet in bullets:
                    p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(18)

        if equation_images:
            # Add equations below bullets
            top = Inches(4) if bullets else Inches(2)

            for img_path in equation_images:
                if img_path and img_path.exists():
                    left = Inches(2)
                    pic = slide.shapes.add_picture(str(img_path), left, top, height=Inches(1))
                    top += Inches(1.2)

    def add_equation_slide(self, title: str, equation_images: List[Path],
                          description: str = ''):
        """Add slide focused on equations."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Blank

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.HEADER_COLOR

        # Add equations centered
        top = Inches(2)
        for img_path in equation_images:
            if img_path and img_path.exists():
                left = Inches(2)
                pic = slide.shapes.add_picture(str(img_path), left, top, height=Inches(1.2))
                top += Inches(1.5)

        # Add description if provided
        if description:
            desc_box = slide.shapes.add_textbox(Inches(1), top + Inches(0.5), Inches(8), Inches(1.5))
            desc_frame = desc_box.text_frame
            desc_frame.text = description
            desc_frame.paragraphs[0].font.size = Pt(16)

    def save(self, output_file: Path):
        """Save presentation."""
        self.prs.save(str(output_file))


def convert_chapter_to_slides(qmd_file: str, output_dir: Optional[str] = None,
                              template: Optional[str] = None):
    """Main conversion function."""

    # Setup paths
    qmd_path = Path(qmd_file)
    if output_dir:
        out_dir = Path(output_dir)
    else:
        out_dir = qmd_path.parent

    chapter_name = qmd_path.stem
    slides_name = f"Slides_{chapter_name.replace('Chapter_', '')}"

    equations_dir = out_dir / f"{slides_name}_equations"
    figures_dir = out_dir / f"{slides_name}_figures"

    # Parse chapter
    print(f"Parsing {qmd_file}...")
    parser = ChapterParser(qmd_file)

    title = parser.extract_title()
    sections = parser.extract_sections()
    equations = parser.extract_equations()

    print(f"  Title: {title}")
    print(f"  Sections: {len(sections)}")
    print(f"  Equations: {len(equations)}")

    # Render equations
    print("Rendering equations...")
    renderer = EquationRenderer(equations_dir, preamble=parser.macros)
    equation_images = renderer.render_all(equations)
    print(f"  Rendered {len(equation_images)}/{len(equations)} equations")

    # Build presentation
    print("Building presentation...")
    builder = SlideBuilder(template)

    # Title slide
    builder.add_title_slide(title)

    # Process sections
    eq_index = 0
    for section in sections:
        # Section title slide
        builder.add_section_slide(section['title'])

        # Extract bullets from section content
        bullets = []
        section_equations = []

        # Simple bullet extraction (lines starting with -, *, or numbered)
        for line in section['content'].split('\n'):
            line = line.strip()
            if line.startswith(('-', '*', '1.', '2.', '3.', '4.', '5.')):
                bullets.append(line.lstrip('-*0123456789. '))

        # Find equations in this section
        section_start = section['content'][:100]  # Use first 100 chars as marker
        for i, eq in enumerate(equations):
            if eq['position'] > parser.content.find(section_start):
                if eq_index < len(equation_images):
                    section_equations.append(equation_images[eq_index])
                    eq_index += 1
                if len(section_equations) >= 3:  # Limit equations per section
                    break

        # Add content slide
        if bullets or section_equations:
            builder.add_content_slide(section['title'], bullets[:5], section_equations[:2])

    # Save presentation
    output_file = out_dir / f"{slides_name}.pptx"
    builder.save(output_file)

    print(f"\nCreated: {output_file}")
    print(f"Equations: {equations_dir}")

    return output_file


def main():
    parser = argparse.ArgumentParser(description='Convert QMD chapter to PowerPoint slides')
    parser.add_argument('qmd_file', help='Chapter QMD file')
    parser.add_argument('--output-dir', '-o', help='Output directory')
    parser.add_argument('--template', '-t', help='PowerPoint template file')

    args = parser.parse_args()

    if not Path(args.qmd_file).exists():
        print(f"Error: File not found: {args.qmd_file}")
        sys.exit(1)

    convert_chapter_to_slides(args.qmd_file, args.output_dir, args.template)


if __name__ == '__main__':
    main()
